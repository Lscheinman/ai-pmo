# backend/services/report_generator.py
from __future__ import annotations

import io, os, json
from datetime import date, timedelta
from typing import Dict, Any, List
from collections import Counter

from fastapi import HTTPException
from sqlalchemy.orm import Session, selectinload, joinedload

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from urllib.parse import urlparse

from db import models
from services.ai_client import (
    hydrate_graph_node_details, enrich_graph_for_llm, sanitize_graph_for_prompt,
    call_llm, extract_llm_text, normalize_ai_text_and_labels
)

TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "..", "templates")


# ----------------------------- PPT helpers -----------------------------------

def _find_shapes_by_name(slide, name: str):
    return [sh for sh in slide.shapes if getattr(sh, "name", "") == name]

def _replace_textvars_in_shape(shape, mapping: Dict[str, str]):
    """Replace {{VAR}} occurrences in a shape's runs."""
    if not getattr(shape, "has_text_frame", False) or not shape.text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            s = r.text or ""
            for k, v in mapping.items():
                s = s.replace(f"{{{{{k}}}}}", v)
            r.text = s

def _set_named_text(slide, shape_name: str, value: str, mapping: Dict[str, str] | None = None):
    """Set text for shapes named `shape_name`; also run {{VAR}} replace across slide if mapping given."""
    for sh in _find_shapes_by_name(slide, shape_name):
        if getattr(sh, "has_text_frame", False):
            sh.text_frame.clear()
            sh.text_frame.paragraphs[0].add_run().text = value
    if mapping:
        for sh in slide.shapes:
            _replace_textvars_in_shape(sh, mapping)

def _replace_text_anywhere(slide, replacements: Dict[str, str]):
    """Replace substring tokens in all runs across the slide."""
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False) or not sh.text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                txt = r.text or ""
                for k, v in replacements.items():
                    if k in txt:
                        txt = txt.replace(k, v)
                r.text = txt

def _add_table_at_anchor(slide, anchor_name: str, rows: int, cols: int, col_widths: List[float] | None = None):
    """Replace a named rectangle anchor with a real table."""
    anchors = _find_shapes_by_name(slide, anchor_name)
    if not anchors:
        return None
    anchor = anchors[0]
    left, top, width, height = anchor.left, anchor.top, anchor.width, anchor.height
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w_in in enumerate(col_widths[:cols]):
            tbl.columns[i].width = Inches(w_in)
    # remove anchor
    try:
        slide.shapes._spTree.remove(anchor._element)
    except Exception:
        pass
    return tbl

def _render_link_chips(slide, links: List[Dict[str, str]]):
    """
    Render "chips" for project links starting at LINK_CHIPS_ANCHOR if present.
    Falls back to a sensible position if no anchor exists.
    """
    anchor = None
    for sh in slide.shapes:
        if getattr(sh, "name", "") == "LINK_CHIPS_ANCHOR":
            anchor = sh
            break

    if anchor:
        left, top = anchor.left, anchor.top
        cols = 3
        try: slide.shapes._spTree.remove(anchor._element)
        except Exception: pass
    else:
        left, top = Inches(6.0), Inches(1.2)
        cols = 2

    chip_w, chip_h = Inches(2.2), Inches(0.5)
    pad_x, pad_y = Inches(0.2), Inches(0.15)

    for idx, link in enumerate(links):
        row, col = divmod(idx, cols)
        x = left + col * (chip_w + pad_x)
        y = top  + row * (chip_h + pad_y)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, chip_w, chip_h)

        # simple style
        fill = shape.fill; fill.solid(); fill.fore_color.rgb = RGBColor(230, 230, 230)
        line = shape.line; line.color.rgb = RGBColor(180, 180, 180)

        if shape.text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = (link.get("title") or "Link")[:40]
            run.font.size = Pt(12)

        try:
            shape.click_action.hyperlink.address = link.get("url")
        except Exception:
            pass


# -------------------------- Data/rollups (PII safe) --------------------------

def _person_display(p: models.Person) -> str:
    """Prefer name; fallback to email; else token."""
    if p and p.name and p.name.strip():
        return p.name.strip()
    if p and p.email and p.email.strip():
        return p.email.strip()
    return f"person_{getattr(p,'id',0)}"

def _raci_people_for_project(db: Session, project_id: int) -> dict:
    """
    Collect R/A/C/I people from ProjectLead.role and TaskAssignee.role.
    Dedup by person_id, stable order (leads first).
    """
    p: models.Project = (
        db.query(models.Project)
          .options(
              selectinload(models.Project.project_leads).selectinload(models.ProjectLead.person),
              selectinload(models.Project.tasks)
                  .selectinload(models.Task.task_assignees).selectinload(models.TaskAssignee.person),
          )
          .filter(models.Project.id == project_id)
          .first()
    )
    out = {"Responsible": [], "Accountable": [], "Consulted": [], "Informed": []}
    seen = set()

    def add(role_label: str, person: models.Person):
        if not person: return
        if person.id in seen: return
        seen.add(person.id)
        out[role_label].append(person)

    for pl in (p.project_leads or []):
        r = (pl.role or "").strip().capitalize()
        if r in out: add(r, pl.person)

    for t in (p.tasks or []):
        for a in (t.task_assignees or []):
            r = (a.role or "").strip().capitalize()
            if r in out: add(r, a.person)

    return out

def _select_project_links(db: Session, project_id: int, limit: int | None = None) -> List[dict]:
    rows = (
        db.query(models.ProjectLink)
          .filter(models.ProjectLink.project_id == project_id)
          .order_by(models.ProjectLink.is_pinned.desc(),
                    models.ProjectLink.sort_order.asc(),
                    models.ProjectLink.created_at.desc())
          .all()
    )
    out = []
    for x in rows[: (limit or len(rows))]:
        title = (x.title or "").strip()
        if not title:
            try:
                title = urlparse(x.url).netloc
            except Exception:
                title = "Link"
        out.append({"title": title, "url": x.url})
    return out

def _task_one_liner_comment(r: dict) -> str:
    """PII-free comment: due window + blocked flag + generic nudge."""
    due = r["due"]
    status = (r["status"] or "").lower()
    tip = []
    if due and due != "—": tip.append(f"Due {due}")
    if status == "blocked": tip.append("Unblock first")
    if not tip: tip.append("Monitor progress")
    return " · ".join(tip)

def _collect_project_bundle(db: Session, project_id: int) -> Dict[str, Any]:
    p: models.Project = (
        db.query(models.Project)
        .options(
            selectinload(models.Project.tags),
            selectinload(models.Project.project_leads).selectinload(models.ProjectLead.person),
            selectinload(models.Project.tasks)
                .selectinload(models.Task.task_assignees).selectinload(models.TaskAssignee.person),
            selectinload(models.Project.tasks).selectinload(models.Task.tags),
            selectinload(models.Project.tasks).selectinload(models.Task.checklist_items),
        )
        .filter(models.Project.id == project_id, models.Project.is_archived == False)  # noqa
        .first()
    )
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")

    task_counts = Counter((t.status or "not started").lower() for t in p.tasks or [])
    total_tasks = len(p.tasks or [])

    # checklist progress
    cl_total = 0
    cl_done = 0
    for t in (p.tasks or []):
        items = t.checklist_items or []
        cl_total += len(items)
        cl_done += sum(1 for ci in items if (ci.status or "").lower() == "complete")

    # due windows
    today = date.today()
    soon7, soon30 = today + timedelta(days=7), today + timedelta(days=30)
    due_stats = {
        "overdue": sum(1 for t in p.tasks or [] if t.end and t.end < today and (t.status or "").lower() != "complete"),
        "due_7":   sum(1 for t in p.tasks or [] if t.end and today <= t.end <= soon7),
        "due_30":  sum(1 for t in p.tasks or [] if t.end and today <= t.end <= soon30),
    }

    tag_counter = Counter()
    for t in p.tasks or []:
        for tg in (t.tags or []):
            if tg and tg.name:
                tag_counter[tg.name] += 1
    top_tags = tag_counter.most_common(10)

    rows = []
    for t in sorted(p.tasks or [], key=lambda x: (x.end or date.max)):
        assignee_ids = [f"person_{a.person_id}" for a in (t.task_assignees or []) if a.person_id]
        rows.append({
            "task": t.name or f"task_{t.id}",
            "task_id": t.id,
            "status": (t.status or "not started"),
            "due": (t.end.isoformat() if t.end else "—"),
            "assignees": ", ".join(assignee_ids) if assignee_ids else "—",
            "tags": ", ".join(sorted({tg.name for tg in (t.tags or []) if tg and tg.name})) or "—",
            "priority": (t.priority or "medium"),
            "type": (t.type or "")  # optional
        })

    return {
        "project": {
            "id": p.id,
            "name": p.name,
            "status": p.status,
            "start": p.start_date.isoformat() if p.start_date else None,
            "end": p.end_date.isoformat() if p.end_date else None,
            "tags": [tg.name for tg in (p.tags or []) if tg and tg.name],
            "description": (p.description or "").strip(),
        },
        "counts": {
            "total_tasks": total_tasks,
            "by_status": dict(task_counts),
            "checklist": {"done": cl_done, "total": cl_total},
            "due": due_stats,
        },
        "top_tags": top_tags,
        "rows": rows,
    }


# ---------------------------- AI (sanitized) ---------------------------------

def _ai_summary_for_project(db: Session, project_id: int) -> str:
    """Sanitized 5–8 bullet summary grounded in IDs/tags only."""
    bundle = _collect_project_bundle(db, project_id)

    nodes = [{
        "data": {
            "id": f"project_{bundle['project']['id']}",
            "type": "Project",
            "label": bundle["project"]["name"],
            "status": bundle["project"]["status"],
            "detail": {
                "description": bundle["project"]["description"][:320],
                "start_date": bundle["project"]["start"],
                "end_date": bundle["project"]["end"],
            },
            "tags_inline": bundle["project"]["tags"][:12]
        }
    }]
    edges = []
    for r in bundle["rows"][:200]:
        tid = f"task_{r['task_id']}"
        nodes.append({"data": {
            "id": tid, "type": "Task", "label": r["task"], "status": r["status"],
            "detail": {"priority": r["priority"], "description_snippet": None},
            "tags_inline": [t.strip() for t in (r["tags"].split(",") if r["tags"] and r["tags"] != "—" else [])][:12]
        }})
        edges.append({"data": {"source": f"project_{bundle['project']['id']}", "target": tid, "type": "PROJECT_TASK"}})
        if r["assignees"] and r["assignees"] != "—":
            for pid in [x.strip() for x in r["assignees"].split(",")]:
                nodes.append({"data": {"id": pid, "type": "Person", "label": pid}})
                edges.append({"data": {"source": tid, "target": pid, "type": "TASK_ASSIGNEE"}})

    graph = {"nodes": nodes, "edges": edges}
    graph = hydrate_graph_node_details(db, graph)
    n, e = enrich_graph_for_llm(graph["nodes"], graph["edges"], db)
    safe = sanitize_graph_for_prompt(n, e)

    prompt = f"""
You are a PMO analyst. Using ONLY the sanitized graph (IDs, status, tags, snippets),
write a crisp project status summary with:
- 5–8 bullets
- progress signals (counts, due windows), risks/blockers, and next steps
- reference entities ONLY by IDs like [task_12] or [person_7]
- do NOT include any names or emails

Graph nodes (subset):
{json.dumps(safe['nodes'][:120], ensure_ascii=False)}

Graph edges (subset):
{json.dumps(safe['edges'][:240], ensure_ascii=False)}
""".strip()

    messages = [{"role": "user", "content": [{"type": "text", "text": prompt}]}]
    resp = call_llm(model_name="gpt-4o", messages=messages, prompt_type="report.summary",
                    meta={"project_id": project_id}, redact_for_log=False)
    raw = extract_llm_text(resp)
    summary, _ = normalize_ai_text_and_labels(raw, db)
    return summary.strip()


# ---------------------------- Slide fillers ----------------------------------

def _fill_slide_1_header(slide, bundle: Dict[str, Any], summary_text: str):
    mapping = {
        "TITLE": f"{bundle['project']['name']} — Summary Report",
        "PROJECT_NAME": bundle['project']['name'],
        "PROJECT_STATUS": bundle['project']['status'] or "—",
        "DATE": date.today().isoformat(),
        "DATE_RANGE": f"{bundle['project']['start'] or '—'} – {bundle['project']['end'] or '—'}",
        "TAGS": (", ".join(bundle['project']['tags']) or "—"),
        "TASK_COUNTS": (
            f"Total {bundle['counts']['total_tasks']} | "
            f"not started {bundle['counts']['by_status'].get('not started',0)} | "
            f"started {bundle['counts']['by_status'].get('started',0)} | "
            f"blocked {bundle['counts']['by_status'].get('blocked',0)} | "
            f"complete {bundle['counts']['by_status'].get('complete',0)}"
        ),
        "CHECKLIST_PROGRESS": (
            "—" if bundle["counts"]["checklist"]["total"] == 0
            else f"{bundle['counts']['checklist']['done']}/{bundle['counts']['checklist']['total']} "
                 f"({round(100*bundle['counts']['checklist']['done']/max(1,bundle['counts']['checklist']['total']))}%)"
        ),
        "DUE_WINDOWS": f"Overdue {bundle['counts']['due']['overdue']} | ≤7d {bundle['counts']['due']['due_7']} | ≤30d {bundle['counts']['due']['due_30']}",
        "OVERVIEW": (bundle['project']['description'][:2000] if bundle['project']['description'] else "—"),
        "SUMMARY_UPDATE": summary_text[:4000],
    }
    for k, v in mapping.items():
        _set_named_text(slide, k, v, mapping)

def _fill_slide_2_agenda(slide, db: Session, project: models.Project):
    raci = _raci_people_for_project(db, project.id)
    resp = _person_display(raci["Responsible"][0]) if raci["Responsible"] else "—"
    acct = _person_display(raci["Accountable"][0]) if raci["Accountable"] else "—"

    title = f"{project.name} — {project.status or '—'}"
    _set_named_text(slide, "AGENDA_TITLE", title)
    _set_named_text(slide, "AGENDA_RESPONSIBLE", f"Responsible: {resp}")
    _set_named_text(slide, "AGENDA_ACCOUNTABLE", f"Accountable: {acct}")

    # also support inline token replacement (if your template uses them)
    _replace_text_anywhere(slide, {
        "PROJECT_STATUS_REPORT": title,
        "RESPONSIBLE": f"Responsible: {resp}",
        "ACCOUNTABLE": f"Accountable: {acct}",
    })

def _fill_slide_3_status(slide, db: Session, project: models.Project, bundle: Dict[str, Any]):
    # RACI token expansion (R0/A0/C0/I0…)
    raci = _raci_people_for_project(db, project.id)
    raci_lists = {
        "R": [ _person_display(p) for p in raci["Responsible"] ],
        "A": [ _person_display(p) for p in raci["Accountable"] ],
        "C": [ _person_display(p) for p in raci["Consulted"] ],
        "I": [ _person_display(p) for p in raci["Informed"] ],
    }
    repl = {}
    for prefix, arr in raci_lists.items():
        for idx, val in enumerate(arr):
            repl[f"{prefix}{idx}"] = val
    if repl:
        _replace_text_anywhere(slide, repl)

    # Optional full RACI table
    max_rows = max([len(x) for x in raci_lists.values()] + [0])
    if max_rows > 0:
        raci_tbl = _add_table_at_anchor(slide, "RACI_TABLE", rows=1 + max_rows, cols=4,
                                        col_widths=[2.2, 2.2, 2.2, 2.2])
        if raci_tbl:
            headers = ["Responsible", "Accountable", "Consulted", "Informed"]
            for j, h in enumerate(headers):
                raci_tbl.cell(0, j).text = h
                raci_tbl.cell(0, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for col, key in enumerate(["R", "A", "C", "I"]):
                arr = raci_lists[key]
                for i in range(max_rows):
                    raci_tbl.cell(i+1, col).text = arr[i] if i < len(arr) else ""

    # Tasks table (all tasks)
    headers = ["Task", "Type", "Priority", "Status", "Start", "End", "Assignees", "Comments"]
    tbl = _add_table_at_anchor(slide, "TASKS_TABLE", rows=1 + len(bundle["rows"]), cols=len(headers),
                               col_widths=[3.5, 1.1, 1.0, 1.2, 1.2, 1.2, 1.8, 3.2])
    if tbl:
        for j, h in enumerate(headers):
            tbl.cell(0, j).text = h
            tbl.cell(0, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for i, r in enumerate(bundle["rows"], start=1):
            assignees_str = r["assignees"]
            if assignees_str and assignees_str != "—":
                ids = [x.strip() for x in assignees_str.split(",") if x.strip()]
                assignees_out = f"{len(ids)} ({', '.join(ids[:6])}{'…' if len(ids) > 6 else ''})"
            else:
                assignees_out = "0"

            vals = [
                r["task"], r.get("type") or "", r["priority"], r["status"],
                bundle["project"]["start"] or "", r["due"],
                assignees_out, _task_one_liner_comment(r)
            ]
            for j, v in enumerate(vals):
                tbl.cell(i, j).text = v

    # Links as chips
    links = _select_project_links(db, project.id)
    if links:
        _render_link_chips(slide, links)


# ---------------------------- Public builder ---------------------------------

def build_project_summary_pptx(db: Session, project_id: int, template_filename: str = "summary_report.pptx") -> bytes:
    """Main entry: returns PPTX bytes for the project summary report."""
    template_path = os.path.join(TEMPLATES_DIR, template_filename)
    if not os.path.exists(template_path):
        raise HTTPException(status_code=404, detail="Template not found")

    bundle = _collect_project_bundle(db, project_id)
    proj = db.query(models.Project).filter(models.Project.id == project_id).first()
    summary_text = _ai_summary_for_project(db, project_id)

    prs = Presentation(template_path)
    if len(prs.slides) >= 1:
        _fill_slide_1_header(prs.slides[0], bundle, summary_text)
    if len(prs.slides) >= 2:
        _fill_slide_2_agenda(prs.slides[1], db, proj)
    if len(prs.slides) >= 3:
        _fill_slide_3_status(prs.slides[2], db, proj, bundle)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


# -------------------------- (Optional) extensibility --------------------------

SUPPORTED_REPORTS = {
    "summary": {
        "template": "summary_report.pptx",
        "builder": build_project_summary_pptx,
    },
    # Add new report types here:
    # "retro": { "template": "retro_report.pptx", "builder": build_project_retro_pptx },
}

def build_report(db: Session, project_id: int, report_type: str = "summary", template_override: str | None = None) -> bytes:
    cfg = SUPPORTED_REPORTS.get(report_type)
    if not cfg:
        raise HTTPException(status_code=400, detail=f"Unsupported report_type '{report_type}'")
    template = template_override or cfg["template"]
    return cfg["builder"](db, project_id, template)
